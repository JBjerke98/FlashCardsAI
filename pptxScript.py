import os
import sys
import json
import fitz  # PyMuPDF for PDF
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI
import re  # Regex for dynamisk sjekk av spørsmålsnummer

# Last inn miljøvariabler fra .env-filen
load_dotenv("/Users/jonasbjerke/Documents/HomeHUB/Egne kode prosjekter/PowerpointFlashcards/api.env")
api_key = os.getenv("OPENAI_API_KEY")

# Initialiser OpenAI-klienten med API-nøkkelen
client = OpenAI(api_key=api_key)

def extract_text_from_pdf(pdf_path):
    """
    Henter ut all tekst fra en PDF-fil.
    """
    pdf_text = ""
    doc = fitz.open(pdf_path)
    
    for page in doc:
        pdf_text += page.get_text()

    doc.close()
    return pdf_text

def extract_text_from_pptx(pptx_path):
    """
    Henter ut all tekst fra en PowerPoint-fil (.pptx).
    """
    presentation = Presentation(pptx_path)
    ppt_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        ppt_text += run.text + " "
    return ppt_text

def parse_questions_and_answers(raw_text):
    """
    Parser råtekst fra OpenAI-responsen til en strukturert JSON-liste.
    """
    qa_pairs = []
    lines = raw_text.split('\n')
    current_question = None
    current_answer = None

    for line in lines:
        line = line.strip()
        if re.match(r"^\d+\.", line):  
            if current_question and current_answer:
                qa_pairs.append({
                    "question": current_question,
                    "answer": current_answer
                })
            # Fjerner nummeret og tar resten som spørsmålet
            current_question = line.split(" ", 1)[1]
            current_answer = ""
        elif line.startswith("-") and current_question:
            current_answer = line.split("-", 1)[1].strip()
        elif current_question:
            current_answer += " " + line

    if current_question and current_answer:
        qa_pairs.append({
            "question": current_question,
            "answer": current_answer
        })

    return qa_pairs

def generate_questions(text):
    """
    Genererer spørsmål basert på teksten ved bruk av OpenAI ChatCompletion.
    """
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages = [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Generate a list of at least 16 well-written questions and answers based on the following text:\n\n{text}. The questions should cover the main topics of the text."}
        ],
        max_tokens=1000
    )
    
    # Hent og parse tekst fra 'message.content'
    raw_text = response.choices[0].message.content
    return json.dumps({"questions_and_answers": parse_questions_and_answers(raw_text)})

def main():
    if len(sys.argv) != 2:
        print("Usage: python script.py <file_path>")
        return

    file_path = sys.argv[1]
    file_extension = os.path.splitext(file_path)[1].lower()

    # Håndterer både PDF- og PPTX-filer
    if file_extension == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif file_extension == ".pdf":
        text = extract_text_from_pdf(file_path)
    else:
        raise ValueError("Unsupported file type. Please provide a .pdf or .pptx file.")

    questions = generate_questions(text)
    print(questions)

if __name__ == "__main__":
    main()
